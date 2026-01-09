"""
PPTX renderer using python-pptx.

Converts SlideElements into editable PowerPoint slides.
"""

import re
from pathlib import Path
from typing import List, Dict, Tuple, Optional
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, PP_PARAGRAPH_ALIGNMENT, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from sliderefactor.models import (
    SlideElements,
    TextBoxElement,
    ImageElement,
    ShapeElement,
    TableElement,
    TableCell,
    BulletItem,
    Slide,
)


def strip_html_tags(text: str) -> str:
    """Strip HTML tags from text and convert common entities."""
    if not text:
        return text
    # Replace <br> tags with newlines
    text = re.sub(r'<br\s*/?>', '\n', text)

    # Replace closing block tags with newlines to prevent word merging
    text = re.sub(r'</p>', '\n', text)
    text = re.sub(r'</div>', '\n', text)
    text = re.sub(r'</h1>', '\n', text)
    text = re.sub(r'</h2>', '\n', text)
    text = re.sub(r'</h3>', '\n', text)

    # Remove all remaining HTML tags
    text = re.sub(r'<[^>]+>', '', text)
    # Convert common HTML entities
    text = text.replace('&amp;', '&')
    text = text.replace('&lt;', '<')
    text = text.replace('&gt;', '>')
    text = text.replace('&nbsp;', ' ')
    text = text.replace('&quot;', '"')
    text = text.replace('&#39;', "'")
    return text.strip()


# Font mapping: PDF internal names -> PowerPoint-compatible names
FONT_MAPPING = {
    # Arial variants
    "arialmt": "Arial",
    "arial-boldmt": "Arial",
    "arial-italicmt": "Arial",
    "arial-bolditalicmt": "Arial",
    "arialmtbold": "Arial",
    "arialnarrow": "Arial Narrow",
    "arialblack": "Arial Black",
    # Times variants
    "timesnewromanpsmt": "Times New Roman",
    "timesnewromanps-boldmt": "Times New Roman",
    "timesnewromanps-italicmt": "Times New Roman",
    "timesnewroman": "Times New Roman",
    "times": "Times New Roman",
    # Helvetica -> Arial (common substitution)
    "helvetica": "Arial",
    "helveticaneue": "Arial",
    "helveticaneue-light": "Arial",
    "helveticaneue-bold": "Arial",
    "helveticaneue-medium": "Arial",
    # Calibri variants
    "calibri": "Calibri",
    "calibri-bold": "Calibri",
    "calibri-italic": "Calibri",
    "calibri-light": "Calibri Light",
    # Cambria
    "cambria": "Cambria",
    "cambriamath": "Cambria Math",
    # Georgia
    "georgia": "Georgia",
    "georgia-bold": "Georgia",
    # Verdana
    "verdana": "Verdana",
    "verdana-bold": "Verdana",
    # Tahoma
    "tahoma": "Tahoma",
    "tahoma-bold": "Tahoma",
    # Trebuchet
    "trebuchetms": "Trebuchet MS",
    "trebuchetms-bold": "Trebuchet MS",
    # Courier/Consolas (monospace)
    "couriernew": "Courier New",
    "couriernewpsmt": "Courier New",
    "courier": "Courier New",
    "consolas": "Consolas",
    # Open Sans
    "opensans": "Open Sans",
    "opensans-regular": "Open Sans",
    "opensans-bold": "Open Sans",
    "opensans-light": "Open Sans",
    # Roboto
    "roboto": "Roboto",
    "roboto-regular": "Roboto",
    "roboto-bold": "Roboto",
    "roboto-light": "Roboto",
    # Segoe UI
    "segoeui": "Segoe UI",
    "segoeui-bold": "Segoe UI",
    "segoeui-light": "Segoe UI Light",
    # Source Sans
    "sourcesanspro": "Source Sans Pro",
    "sourcesanspro-regular": "Source Sans Pro",
    # Lato
    "lato": "Lato",
    "lato-regular": "Lato",
    "lato-bold": "Lato",
}

# Default fallback font
DEFAULT_FONT = "Calibri"


def normalize_font_name(pdf_font_name: Optional[str]) -> str:
    """
    Normalize a PDF font name to a PowerPoint-compatible font name.

    PDF fonts often have internal names like 'ArialMT', 'TimesNewRomanPS-BoldMT',
    'BCDEEE+Calibri', etc. This function maps them to standard names.
    """
    if not pdf_font_name:
        return DEFAULT_FONT

    # Remove common PDF font prefixes (subset prefixes like 'BCDEEE+')
    clean_name = re.sub(r'^[A-Z]{6}\+', '', pdf_font_name)

    # Normalize: lowercase, remove spaces and hyphens for matching
    normalized = clean_name.lower().replace(' ', '').replace('-', '').replace('_', '')

    # Check direct mapping
    if normalized in FONT_MAPPING:
        return FONT_MAPPING[normalized]

    # Check if any mapping key is contained in the normalized name
    for key, value in FONT_MAPPING.items():
        if key in normalized:
            return value

    # Check for common font family names
    common_families = [
        ("arial", "Arial"),
        ("helvetica", "Arial"),
        ("times", "Times New Roman"),
        ("calibri", "Calibri"),
        ("cambria", "Cambria"),
        ("georgia", "Georgia"),
        ("verdana", "Verdana"),
        ("tahoma", "Tahoma"),
        ("trebuchet", "Trebuchet MS"),
        ("courier", "Courier New"),
        ("consolas", "Consolas"),
        ("segoe", "Segoe UI"),
        ("roboto", "Roboto"),
        ("opensans", "Open Sans"),
        ("lato", "Lato"),
    ]

    for pattern, replacement in common_families:
        if pattern in normalized:
            return replacement

    # If no match found, return the original (cleaned) name
    # PowerPoint will use a fallback if it doesn't recognize the font
    return clean_name if clean_name else DEFAULT_FONT


class PPTXRenderer:
    """
    Render SlideElements into a PowerPoint presentation using python-pptx.

    Features:
    - Text boxes with editable content
    - Bullet lists with proper nesting
    - Images with positioning
    - Style preservation (alignment, fonts, etc.)
    """

    # Style size mapping (relative to slide height)
    SIZE_MAP = {
        "xs": 0.015, # 1.5% (approx 8pt)
        "sm": 0.02,  # 2% (approx 11pt)
        "md": 0.03,  # 3% (approx 16pt)
        "lg": 0.045, # 4.5% (approx 24pt)
        "xl": 0.06,  # 6% (approx 32pt)
    }

    # Role-based default sizes
    ROLE_SIZE_MAP = {
        "title": "xl",
        "subtitle": "lg",
        "body": "md",
        "caption": "sm",
        "footer": "xs",
    }

    def __init__(
        self,
        slide_width_inches: float = 10.0,
        slide_height_inches: float = 7.5,
        dpi: int = 96,
        render_background: bool = True,
    ):
        """
        Initialize renderer.

        Args:
            slide_width_inches: Slide width in inches
            slide_height_inches: Slide height in inches
            dpi: DPI for pixel-to-inch conversion
            render_background: Whether to render background images (disable to avoid "double text")
        """
        self.slide_width_inches = slide_width_inches
        self.slide_height_inches = slide_height_inches
        self.dpi = dpi
        self.render_background = render_background

    def render(
        self,
        elements_list: List[SlideElements],
        slides_info: List[Slide],
        output_path: Path,
        images_dir: Path,
    ) -> Path:
        """
        Render all slides to a PPTX file.

        Args:
            elements_list: List of SlideElements for each slide
            slides_info: List of original Slide objects (for dimensions)
            output_path: Path to save the PPTX file
            images_dir: Directory containing extracted images

        Returns:
            Path to the generated PPTX file
        """
        prs = Presentation()

        # Calculate slide dimensions from first slide's aspect ratio
        if slides_info:
            first_slide = slides_info[0]
            aspect_ratio = first_slide.width_px / first_slide.height_px
            
            # Use standard height and calculate width to match aspect ratio
            # Common approach: 7.5" height for widescreen presentations
            self.slide_height_inches = 7.5
            self.slide_width_inches = self.slide_height_inches * aspect_ratio
            
            print(f"[PPTX] Source aspect ratio: {aspect_ratio:.3f} ({first_slide.width_px}x{first_slide.height_px}px)")
            print(f"[PPTX] Slide dimensions: {self.slide_width_inches:.2f}\" x {self.slide_height_inches:.2f}\"")

        # Set slide dimensions
        prs.slide_width = Inches(self.slide_width_inches)
        prs.slide_height = Inches(self.slide_height_inches)

        print(f"[PPTX] Rendering {len(elements_list)} slides")

        for i, (elements, slide_info) in enumerate(zip(elements_list, slides_info)):
            print(
                f"[PPTX] Rendering slide {i + 1}/{len(elements_list)} ({len(elements.elements)} elements)"
            )

            # Use blank slide layout
            slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)

            # Render background if present and enabled
            if self.render_background and slide_info.background.mode == "image" and slide_info.background.image_ref:
                self._render_background(slide, slide_info, images_dir)

            # Render each element
            for element in elements.elements:
                if isinstance(element, TextBoxElement):
                    self._render_textbox(element, slide, slide_info)
                elif isinstance(element, ImageElement):
                    self._render_image(element, slide, slide_info, images_dir)
                elif isinstance(element, ShapeElement):
                    self._render_shape(element, slide, slide_info)
                elif isinstance(element, TableElement):
                    self._render_table(element, slide, slide_info)

        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))

        print(f"[PPTX] Saved presentation to {output_path}")
        return output_path

    def _render_textbox(
        self, element: TextBoxElement, slide, slide_info: Slide
    ) -> None:
        """Render a text box element."""
        # Convert bbox from pixels to inches
        print(f"[PPTX] Debug: Slide px={slide_info.width_px}x{slide_info.height_px}, BBox={element.bbox.to_list()[:2]}")
        left, top, width, height = self._bbox_to_inches(
            element.bbox.to_list(),
            slide_info.width_px,
            slide_info.height_px,
        )

        # Add text box
        # Slightly widen the box to prevent premature wrapping
        width += 0.2
        textbox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        
        # Minimize margins for tighter layout
        text_frame.margin_left = 0
        text_frame.margin_right = 0
        text_frame.margin_top = 0
        text_frame.margin_bottom = 0

        # Set vertical alignment
        if element.style_hints.vertical_align == "middle":
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        elif element.style_hints.vertical_align == "bottom":
            text_frame.vertical_anchor = MSO_ANCHOR.BOTTOM
        else:
            text_frame.vertical_anchor = MSO_ANCHOR.TOP

        # Determine font size
        # Use output slide height in POINTS (not pixels) for proper font sizing
        # 1 inch = 72 points
        slide_height_pt = self.slide_height_inches * 72
        size_key = element.style_hints.size or self.ROLE_SIZE_MAP.get(element.role, "md")
        font_size = int(self.SIZE_MAP[size_key] * slide_height_pt)

        # Add content
        if element.structure.type == "bullets":
            self._add_bullets(
                text_frame,
                element.structure.items,
                element.style_hints,
                font_size,
                element.font_hints,
            )
        else:  # paragraphs
            self._add_paragraphs(
                text_frame,
                element.structure.items,
                element.style_hints,
                font_size,
                element.font_hints,
            )

    def _add_bullets(
        self,
        text_frame,
        items: List[BulletItem],
        style_hints,
        font_size: int,
        font_hints,
    ) -> None:
        """Add bullet points to text frame."""
        text_frame.clear()  # Remove default paragraph

        for i, item in enumerate(items):
            if isinstance(item, str):
                # Simple string bullet
                p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
                p.text = strip_html_tags(item)
                p.level = 0
            elif isinstance(item, BulletItem):
                p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
                p.level = item.level

                # Add runs
                if item.runs:
                    for run_data in item.runs:
                        run = p.add_run()
                        run.text = strip_html_tags(run_data.text)

                        if run_data.bold:
                            run.font.bold = True
                        if run_data.italic:
                            run.font.italic = True
                        if run_data.underline:
                            run.font.underline = True
                        if run_data.font_size:
                            run.font.size = Pt(run_data.font_size)
                        elif font_hints and font_hints.size:
                            run.font.size = Pt(font_hints.size)
                        else:
                            run.font.size = Pt(font_size)
                        if run_data.font_name:
                            run.font.name = normalize_font_name(run_data.font_name)
                        elif font_hints and font_hints.name:
                            run.font.name = normalize_font_name(font_hints.name)
                else:
                    # No runs, just set text
                    p.text = strip_html_tags(item.text)

            # Apply paragraph-level styling
            if style_hints.align == "center":
                p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            elif style_hints.align == "right":
                p.alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT
            else:
                p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

            # Set font size for the paragraph
            if not item.runs if isinstance(item, BulletItem) else True:
                for run in p.runs:
                    if run.font.size is None:
                        if font_hints and font_hints.size:
                            run.font.size = Pt(font_hints.size)
                        else:
                            run.font.size = Pt(font_size)
                    if font_hints and font_hints.name:
                        run.font.name = normalize_font_name(font_hints.name)

            # Apply weight
            if style_hints.weight == "bold":
                for run in p.runs:
                    run.font.bold = True

    def _add_paragraphs(
        self,
        text_frame,
        items: List[str],
        style_hints,
        font_size: int,
        font_hints,
    ) -> None:
        """Add paragraphs to text frame."""
        text_frame.clear()

        for i, text in enumerate(items):
            p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
            # Strip HTML tags from text
            p.text = strip_html_tags(str(text))

            # Apply styling
            if style_hints.align == "center":
                p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            elif style_hints.align == "right":
                p.alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT
            else:
                p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

            for run in p.runs:
                # Apply font size
                if font_hints and font_hints.size:
                    run.font.size = Pt(font_hints.size)
                else:
                    run.font.size = Pt(font_size)
                
                # Apply font name
                if font_hints and font_hints.name:
                    run.font.name = normalize_font_name(font_hints.name)
                
                # Apply bold (from font_hints or style_hints)
                if (font_hints and font_hints.bold) or style_hints.weight == "bold":
                    run.font.bold = True
                
                # Apply italic
                if font_hints and font_hints.italic:
                    run.font.italic = True
                
                # Apply color
                if font_hints and font_hints.color:
                    color = self._parse_hex_color(font_hints.color)
                    if color:
                        run.font.color.rgb = color

    def _render_image(
        self, element: ImageElement, slide, slide_info: Slide, images_dir: Path
    ) -> None:
        """Render an image element."""
        # Find image file
        image_path = images_dir / element.image_ref

        if not image_path.exists():
            # Fallback: Try to crop from full page image
            page_image_path = images_dir / f"page_{slide_info.page_index}.png"
            if page_image_path.exists():
                try:
                    print(f"[PPTX] Cropping image for {element.image_ref} from {page_image_path}")
                    with Image.open(page_image_path) as img:
                        # Calculate scale factor between bbox coordinates and image pixels
                        # element.bbox is in slide_info.width_px coordinates
                        scale_x = img.width / slide_info.width_px
                        scale_y = img.height / slide_info.height_px
                        
                        x0, y0, x1, y1 = element.bbox.to_list()
                        
                        # Apply scaling to crop box
                        crop_box = (
                            int(x0 * scale_x), 
                            int(y0 * scale_y), 
                            int(x1 * scale_x), 
                            int(y1 * scale_y)
                        )
                        
                        # Ensure crop box is valid
                        if crop_box[2] > crop_box[0] and crop_box[3] > crop_box[1]:
                            cropped_img = img.crop(crop_box)
                            
                            # Create a valid filename
                            safe_name = re.sub(r'[^\w\-]', '_', str(element.image_ref))
                            if not safe_name.lower().endswith(('.png', '.jpg', '.jpeg')):
                                safe_name += ".png"
                                
                            image_path = images_dir / safe_name
                            cropped_img.save(image_path)
                            print(f"[PPTX] Saved cropped image to {image_path}")
                        else:
                            print(f"[PPTX] Warning: Invalid crop box {crop_box}")
                            return
                except Exception as e:
                    print(f"[PPTX] Error cropping image: {e}")
                    return
            else:
                print(f"[PPTX] Warning: Image not found: {image_path} and page image missing")
                return

        # Convert bbox from pixels to inches
        left, top, width, height = self._bbox_to_inches(
            element.bbox.to_list(),
            slide_info.width_px,
            slide_info.height_px,
        )

        try:
            # Add picture
            slide.shapes.add_picture(
                str(image_path),
                Inches(left),
                Inches(top),
                width=Inches(width),
                height=Inches(height),
            )
        except Exception as e:
            print(f"[PPTX] Warning: Failed to add image {image_path}: {e}")

    def _render_background(self, slide, slide_info: Slide, images_dir: Path) -> None:
        """Render slide background image."""
        image_path = images_dir / slide_info.background.image_ref
        if not image_path.exists():
            print(f"[PPTX] Warning: Background image not found: {image_path}")
            return

        try:
            slide.shapes.add_picture(
                str(image_path),
                Inches(0),
                Inches(0),
                width=Inches(self.slide_width_inches),
                height=Inches(self.slide_height_inches),
            )
        except Exception as e:
            print(f"[PPTX] Warning: Failed to add background {image_path}: {e}")

    def _render_shape(
        self, element: ShapeElement, slide, slide_info: Slide
    ) -> None:
        """Render a shape element."""
        # Convert bbox from pixels to inches
        left, top, width, height = self._bbox_to_inches(
            element.bbox.to_list(),
            slide_info.width_px,
            slide_info.height_px,
        )

        # Map shape type
        shape_type_map = {
            "rectangle": MSO_SHAPE.RECTANGLE,
            "circle": MSO_SHAPE.OVAL,
            "line": MSO_SHAPE.ROUNDED_RECTANGLE,  # Placeholder
            "arrow": MSO_SHAPE.RIGHT_ARROW,
        }

        shape_type = shape_type_map.get(element.shape_type, MSO_SHAPE.RECTANGLE)

        # Add shape
        shape = slide.shapes.add_shape(
            shape_type,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height),
        )

        # Apply fill color
        if element.fill_color:
            color = self._parse_hex_color(element.fill_color)
            if color:
                shape.fill.solid()
                shape.fill.fore_color.rgb = color

        # Apply border
        if element.border_color:
            color = self._parse_hex_color(element.border_color)
            if color:
                shape.line.color.rgb = color
                shape.line.width = Pt(element.border_width)

    def _bbox_to_inches(
        self, bbox: List[float], slide_width_px: float, slide_height_px: float
    ) -> Tuple[float, float, float, float]:
        """
        Convert bounding box from pixels to inches.

        Args:
            bbox: [x0, y0, x1, y1] in pixels
            slide_width_px: Slide width in pixels
            slide_height_px: Slide height in pixels

        Returns:
            (left, top, width, height) in inches
        """
        x0, y0, x1, y1 = bbox

        # Convert to relative coordinates (0-1)
        x0_rel = x0 / slide_width_px
        y0_rel = y0 / slide_height_px
        x1_rel = x1 / slide_width_px
        y1_rel = y1 / slide_height_px

        # Convert to inches
        left = x0_rel * self.slide_width_inches
        top = y0_rel * self.slide_height_inches
        width = (x1_rel - x0_rel) * self.slide_width_inches
        height = (y1_rel - y0_rel) * self.slide_height_inches

        return left, top, width, height

    @staticmethod
    def _parse_hex_color(hex_color: str) -> Optional[RGBColor]:
        """Parse hex color string to RGBColor."""
        try:
            hex_color = hex_color.lstrip("#")
            if len(hex_color) == 6:
                r = int(hex_color[0:2], 16)
                g = int(hex_color[2:4], 16)
                b = int(hex_color[4:6], 16)
                return RGBColor(r, g, b)
        except ValueError:
            pass
        return None
